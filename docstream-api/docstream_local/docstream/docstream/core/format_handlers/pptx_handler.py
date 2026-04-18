"""
PPTX Handler — extracts content from PowerPoint files.
Uses python-pptx to extract slide titles, body text,
speaker notes, and table content.
"""

from __future__ import annotations

from pathlib import Path

from docstream.exceptions import ExtractionError
from docstream.models.document import Block, BlockType


def _table_to_markdown(table) -> str:  # type: ignore[type-arg]
    """Convert a python-pptx Table object to a Markdown table string."""
    rows: list[list[str]] = []
    for row in table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    if not rows:
        return ""
    header = rows[0]
    separator = ["---"] * len(header)
    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in rows[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


class PPTXHandler:
    """Extract blocks from ``.pptx`` PowerPoint presentations.

    Processes each slide in order, extracting the title,
    body text shapes, speaker notes, and any table shapes
    using ``python-pptx``.
    """

    def extract(self, file_path: Path) -> list[Block]:
        """Extract blocks from a ``.pptx`` file.

        Args:
            file_path: Path to the ``.pptx`` file.

        Returns:
            List of ``Block`` objects — one per slide element
            (title, body, notes, table rows).

        Raises:
            ExtractionError: If the file cannot be opened or parsed.
        """
        try:
            from pptx import Presentation  # python-pptx
        except ImportError as exc:
            raise ExtractionError(
                "python-pptx is required for PPTX extraction. "
                "Install with: pip install python-pptx"
            ) from exc

        try:
            prs = Presentation(str(file_path))
        except Exception as exc:
            raise ExtractionError("Could not open PPTX file.") from exc

        blocks: list[Block] = []

        for slide_index, slide in enumerate(prs.slides):
            page_num = slide_index + 1

            # 1. Slide title
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame:
                title_text = title_shape.text.strip()
                if title_text:
                    blocks.append(
                        Block(
                            type=BlockType.HEADING,
                            content=title_text,
                            page_number=page_num,
                            font_size=28.0,
                            is_bold=True,
                            bbox=(0.0, 0.0, 0.0, 0.0),
                        )
                    )

            # 2. Body text frames + tables
            for shape in slide.shapes:
                # Skip title (already handled)
                if shape == title_shape:
                    continue

                # Tables
                if shape.has_table:
                    markdown = _table_to_markdown(shape.table)
                    if markdown:
                        blocks.append(
                            Block(
                                type=BlockType.TABLE,
                                content=markdown,
                                page_number=page_num,
                                font_size=12.0,
                                bbox=(0.0, 0.0, 0.0, 0.0),
                            )
                        )
                    continue

                # Text frames
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        content = para.text.strip()
                        if not content:
                            continue
                        blocks.append(
                            Block(
                                type=BlockType.TEXT,
                                content=content,
                                page_number=page_num,
                                font_size=12.0,
                                bbox=(0.0, 0.0, 0.0, 0.0),
                            )
                        )

            # 3. Speaker notes
            try:
                notes_slide = slide.notes_slide
                if notes_slide:
                    notes_tf = notes_slide.notes_text_frame
                    if notes_tf:
                        notes_text = notes_tf.text.strip()
                        if notes_text:
                            blocks.append(
                                Block(
                                    type=BlockType.TEXT,
                                    content=f"[Speaker notes: {notes_text}]",
                                    page_number=page_num,
                                    font_size=10.0,
                                    bbox=(0.0, 0.0, 0.0, 0.0),
                                )
                            )
            except Exception:
                # Notes are optional — never fail extraction over them
                pass

        return blocks
